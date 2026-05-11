from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()
slide_width = prs.slide_width

# Define a helper function to add a title and content slide
def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = '\n'.join(content_lines)
    return slide

# Define all slide contents based on the outline
slides_data = [
    ("CipherVault – A Secure Password Management System",
     ["Charan Basava", "Course: Computing", "Supervisor: [Your Supervisor’s Name]", "Date: [Your Presentation Date]"]),

    ("Project Context & Rationale", 
     ["Problem: Insecure password storage is a risk.",
      "Solution: CipherVault enables secure, encrypted credential management.",
      "Motivation: Growing cybersecurity threats, usability gaps in existing tools."]),

    ("Aims and Objectives",
     ["Aim: To develop a secure, lightweight password manager.",
      "Objectives:",
      "• Secure login/registration",
      "• Encrypted password storage",
      "• Responsive frontend UI",
      "• User-friendly experience"]),

    ("Functional Requirements (Summary)", 
     ["Examples: Register/Login, Store/Retrieve Credentials",
      "MoSCoW Prioritisation Overview"]),

    ("Non-Functional Requirements",
     ["Examples: Encryption (Fernet), Usability, Maintainability",
      "Highlight Must-Haves and Justifications"]),

    ("High-Level System Architecture", 
     ["Frontend: HTML/CSS/JavaScript",
      "Backend: FastAPI",
      "Database: MySQL + SQLModel",
      "Encryption Layer: Fernet (Python Cryptography)"]),

    ("Key Design Decisions", 
     ["Chose FastAPI for async, performance, Swagger docs",
      "Used Fernet for symmetric encryption",
      "Clean, modular backend for scalability"]),

    ("Models Overview", 
     ["Use Case Diagram: User interactions overview",
      "Class Diagram: User and Credential classes",
      "ERD: Entities and relationships"]),

    ("Testing Strategy", 
     ["Unit testing backend logic",
      "Manual testing of UI + API",
      "Test Results: Login success, Password encryption integrity, Form validation"]),

    ("Demonstration Flow", 
     ["Includes: Registration, Login, Add/View/Delete Passwords",
      "Password strength check feature",
      "Server startup via Tkinter launcher app"]),

    ("Reflection – What Went Well", 
     ["Effective modular design",
      "Achieved all must-have requirements",
      "Smooth frontend-backend integration"]),

    ("Reflection – Challenges", 
     ["Debugging asynchronous FastAPI routes",
      "Encrypting/decrypting user-specific data securely",
      "UI validation edge cases"]),

    ("Lessons Learned", 
     ["Gained strong experience in full-stack security",
      "Handled cryptographic libraries and API routing",
      "Improved testing and debugging strategies"]),

    ("What Would I Do Differently", 
     ["Automate more test cases",
      "Implement full mobile compatibility",
      "Add MFA or biometric login in future iterations"]),

    ("Q&A / Closing Slide", 
     ["Thank you!", "Any Questions?"]),
]

# Add all slides
for title, content in slides_data:
    add_slide(title, content)

# Save presentation
pptx_path = "/mnt/data/CipherVault_Final_Presentation11.pptx"
prs.save(pptx_path)
pptx_path
